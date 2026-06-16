"""
extract.py - Extract plain text from various document formats.
"""

import os


def extract_text(filepath: str) -> str:
    """Extract text from a file based on its extension. Returns '' on failure."""
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".txt" or ext == ".md" or ext == ".csv":
            return _read_plain(filepath)
        elif ext == ".pdf":
            return _read_pdf(filepath)
        elif ext == ".docx":
            return _read_docx(filepath)
        elif ext == ".pptx":
            return _read_pptx(filepath)
        elif ext in (".xlsx", ".xlsm"):
            return _read_xlsx(filepath)
        else:
            return ""
    except Exception as e:
        print(f"  [skip] {filepath}: {e}")
        return ""


def _read_plain(filepath):
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _read_pdf(filepath):
    from pypdf import PdfReader

    reader = PdfReader(filepath)
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        text_parts.append(page_text)
    return "\n".join(text_parts)


def _read_docx(filepath):
    import docx

    doc = docx.Document(filepath)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)

    return "\n".join(parts)


def _read_pptx(filepath):
    from pptx import Presentation

    prs = Presentation(filepath)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
            # speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(notes)
    return "\n".join(parts)


def _read_xlsx(filepath):
    from openpyxl import load_workbook

    wb = load_workbook(filepath, data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


SUPPORTED_EXTENSIONS = {".txt", ".md", ".csv", ".pdf", ".docx", ".pptx", ".xlsx", ".xlsm"}